"""
See http://pbpython.com/creating-powerpoint.html for details on this script
Requires https://python-pptx.readthedocs.org/en/latest/index.html
Program takes a PowerPoint input file and generates a marked up version that
shows the various layouts and placeholders in the template.
"""

from pptx import Presentation
import argparse


def parse_args():
    """ Setup the input and output arguments for the script
    Return the parsed input and output files
    """
    parser = argparse.ArgumentParser(description='Create personalized powerpoint')
    parser.add_argument('--infile',
                        type=argparse.FileType('r'),
                        required=True,
                        dest='inf',
                        help='Powerpoint file to be personalized')
    parser.add_argument('--outfile',
                        type=argparse.FileType('w'),
                        required=True,
                        dest='outf',
                        help='Output powerpoint')
    parser.add_argument('--name',
                        type=str,
                        help='Name to add')  
    parser.add_argument('-a',
                        dest='analyze_flag',
                        default=False,
                        action='store_true',
                        help='analyze ppt only')
    return parser.parse_args()


def analyze_ppt(input, output):
    """ Take the input file and analyze the structure.
    The output file contains marked up information to make it easier
    for generating future powerpoint templates.
    """
    prs = Presentation(input)
    # Each powerpoint file has multiple layouts
    # Loop through them all and  see where the various elements are
    for index, _ in enumerate(prs.slide_layouts):
        slide = prs.slides.add_slide(prs.slide_layouts[index])
        # Not every slide has to have a title
        try:
            title = slide.shapes.title
            title.text = 'Title for Layout {}'.format(index)
        except AttributeError:
            print("No Title for Layout {}".format(index))
        # Go through all the placeholders and identify them by index and type
        for shape in slide.placeholders:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                # Do not overwrite the title which is just a special placeholder
                try:
                    if 'Title' not in shape.text:
                        shape.text = 'Placeholder index:{} type:{}'.format(phf.idx, shape.name)
                except AttributeError:
                    print("{} has no text attribute".format(phf.type))
                print('{} {}'.format(phf.idx, shape.name))
    prs.save(output)

def create_ppt(input, output, name):
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    prs = Presentation(input)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = name 
    prs.save(output)
    

if __name__ == "__main__":
    args = parse_args()
    if(args.analyze_flag):
        analyze_ppt(args.inf.name, args.outf.name)
    else:
        create_ppt(args.inf.name, args.outf.name, args.name)